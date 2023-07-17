import os
import datetime
import argparse
import pandas as pd
from visualize_datasets import (heatmap, boxplot, plot_subset, plot_geomap, 
                    plot_points_for_station, get_first_last_date)

from jinja2 import Environment, FileSystemLoader
#from weasyprint import HTML, CSS

def get_last_commit():
    """ Récupération du dernier commit """
    return os.popen('git log -1 --format="%H"').read().strip()

def create_report_pdf(vars, name, template_path):
    
    report_dir = vars['report_dir']
    env = Environment(loader=FileSystemLoader('.'))
    template = env.get_template(template_path)
    html = template.render(vars)
    #with open(os.path.join(report_dir, name + '.html'), 'w') as f:
    #    f.write(html)
    css = [CSS(string='body { font-family: Arial, Helvetica, sans-serif; }')]
    HTML(string=html).write_pdf(os.path.join(report_dir, name + '.pdf'), stylesheets=css, presentational_hints=True)
    
def create_report_pptx(vars, name):
    from pptx import Presentation
    from pptx.util import Cm, Pt
    #https://buildmedia.readthedocs.org/media/pdf/python-pptx/latest/python-pptx.pdf
    
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6]) # blank slide
    txtBox = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(8), Cm(4))
    tf = txtBox.text_frame
    p = tf.add_paragraph()
    p.font.size = Pt(12)
    p.text = f'''Rapport de données {vars.get("crop", "")} \n\n'''
    for k, v in vars.items():
        if 'graph' in k:
            continue
        p.text += f'{k}: {v}\n'
   
    
    f = vars.get('graph_heatmap', None)
    if f:
        _ = slide.shapes.add_picture(f, Cm(8), Cm(4), Cm(16), Cm(8))
        f.close()
        
    f = vars.get('graph_plot_subset', None)
    if f:
        _ = slide.shapes.add_picture(f, Cm(8), Cm(13), Cm(15), Cm(5))
        f.close()
    
    report_dir = vars['report_dir']
    filename = os.path.join(report_dir, name + '.pptx')
    pres.save(filename)
        
    
if __name__ == '__main__':
    
    parser = argparse.ArgumentParser(prog=os.path.basename(__file__))
    parser.add_argument('-f', '--filepath', help='Path or url of file with data',  action='store', default=os.getcwd())  
    parser.add_argument('-c', '--crop', help='Crop name',  action='store', default=os.getcwd())  
    
    args = parser.parse_args()
    filepath = args.filepath
    crop = args.crop
    filename = os.path.basename(filepath)
    
    today = datetime.datetime.now().strftime("%Y-%m-%d")
    
    report_dir = os.path.join('reports', today)
    if not os.path.exists(report_dir):
        os.makedirs(report_dir)
    
    df = pd.read_csv(filepath)
    
    vars = {'report_dir': report_dir, 
            'git version': get_last_commit(),
            'date': today,
            'crop': crop,
            'filename': filename,
            'n_rows': df.shape[0],
            'n_cols': df.shape[1]
            }
   
    
    #vars['df_head'] = df.head().to_html()   
  
    vars['graph_heatmap'] = heatmap(df, index='Annee', values='Rendement', columns='Region', capture_plot=True)
  
    vars['graph_plot_subset'] =  plot_subset(df, x='Annee', y='Rendement', values=df.Region.unique()[0:4],
                values_col='Region', y_units='boiseaux/acre', capture_plot=True)
   
    create_report_pptx(vars, f'{crop}_{today}')
    
  